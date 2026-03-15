import os
from loguru import logger
from .registry import register_tool
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from datetime import datetime



def _get_default_filename(prefix: str, ext: str) -> str:
    doc_dir = os.path.join("data", "documents")
    os.makedirs(doc_dir, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return os.path.join(doc_dir, f"{prefix}_{timestamp}.{ext}")

@register_tool("create_word_doc")
def create_word_doc(content: str, filepath: str = None) -> str:
    """Creates a Microsoft Word (.docx) document with the given text content."""
    try:
        if not filepath:
            filepath = _get_default_filename("document", "docx")
        
        doc = Document()
        doc.add_paragraph(content)
        doc.save(filepath)
        return f"Word document saved successfully at {os.path.abspath(filepath)}"
    except Exception as e:
        return f"Error creating Word document: {e}"

@register_tool("create_excel_file")
def create_excel_file(sheet_name: str = "Sheet1", headers: list = None, rows: list = None, filepath: str = None) -> str:
    """Creates a Microsoft Excel (.xlsx) file with given headers and rows (list of lists)."""
    try:
        if not filepath:
            filepath = _get_default_filename("spreadsheet", "xlsx")
            
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        
        current_row = 1
        if headers:
            for col_num, header in enumerate(headers, 1):
                ws.cell(row=current_row, column=col_num, value=header)
            current_row += 1
            
        if rows:
            for row_data in rows:
                for col_num, cell_value in enumerate(row_data, 1):
                    ws.cell(row=current_row, column=col_num, value=cell_value)
                current_row += 1
                
        wb.save(filepath)
        return f"Excel workflow saved successfully at {os.path.abspath(filepath)}"
    except Exception as e:
        logger.error(f"Excel error: {e}")
        return f"Error creating Excel document: {e}"

@register_tool("create_powerpoint")
def create_powerpoint(title: str, subtitle: str, bullets: list = None, filepath: str = None) -> str:
    """Creates a basic PowerPoint (.pptx) presentation with a title slide and an optional content slide."""
    try:
        if not filepath:
            filepath = _get_default_filename("presentation", "pptx")
            
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        if bullets:
            # Bullet slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide2 = prs.slides.add_slide(bullet_slide_layout)
            slide2.shapes.title.text = "Details"
            body_shape = slide2.shapes.placeholders[1]
            tf = body_shape.text_frame
            
            for index, bullet in enumerate(bullets):
                p = tf.add_paragraph() if index > 0 else tf.paragraphs[0]
                p.text = str(bullet)
                p.level = 0
                
        prs.save(filepath)
        return f"PowerPoint saved successfully at {os.path.abspath(filepath)}"
    except Exception as e:
        return f"Error creating PowerPoint: {e}"
